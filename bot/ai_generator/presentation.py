import io
import os
import re

from pptx import Presentation

try:
    from image_scrapper import downloader
except ImportError:
    from .image_scrapper import downloader


async def generate_ppt_prompt(language, emotion_type, slide_length, topic):
    message = f"""In your role as a presentation virtuoso, your mission is to meticulously sculpt a {language} language framework for an impactful {emotion_type} slideshow presentation, unfurling the enthralling narrative of {topic} across a robust expanse of {slide_length} slides.
     
Make sure it is {slide_length} slides long.

Your canvas includes a suite of slide types:

Slide Types:
- Title Slide: (Title, Subtitle)
- Content Slide: (Title, Content)
- Image Slide: (Title, Content, Image)
- Thanks Slide: (Title)

Precede each slide type with the prescribed tags:
- Title Slide: [L_TS]
- Content Slide: [L_CS]
- Image Slide: [L_IS]
- Thanks Slide: [L_THS]

Infuse clarity by inserting a [SLIDEBREAK] tag after each slide, ensuring seamless flow and engagement.

Example:
[L_TS]
[TITLE]Mount Everest: The Apex of Achievement[/TITLE]

[SLIDEBREAK]

[L_IS]
[TITLE]Facts about Mount Everest[/TITLE]
[CONTENT]• Towering majestically at an altitude of 8,848 meters (29,029 ft) above sea level
• Conquered triumphantly by Sir Edmund Hillary and Tenzing Norgay on May 29, 1953
• Eclipsed by over 300 lives lost in pursuit of the summit[/CONTENT]
[IMAGE]Mount Everest[/IMAGE]

[SLIDEBREAK]

Embrace the designated tags for consistent formatting:
- Before Title: [TITLE]
- After Title: [/TITLE]
- Before Subtitle: [SUBTITLE]
- After Subtitle: [/SUBTITLE]
- Before Content: [CONTENT]
- After Content: [/CONTENT]
- Before Image: [IMAGE]
- After Image: [/IMAGE]

Elevate each Content segment with comprehensive insights, ensuring an immersive exploration. Conclude with [/CONTENT].

Harmonize language to the presentation's ethos in {language}. Complement visuals with vivid descriptions using evocative keywords such as "Mount Everest Sunset" or "Niagara Falls Rainbow".

Steer clear of referencing "Image" within the Content tag and eschew special characters (?, !, ., :, ) in the Title.

Commit to the prescribed format with unwavering precision, eschewing superfluous embellishments."""

    return message


async def generate_ppt(answer, template):
    template = os.path.join("bot", "ai_generator", "presentation_templates", f"{template}.pptx")
    root = Presentation(template)

    # """ Ref for slide types:
    # 0 -> title and subtitle
    # 1 -> title and content
    # 2 -> section header
    # 3 -> two content
    # 4 -> Comparison
    # 5 -> Title only
    # 6 -> Blank
    # 7 -> Content with caption
    # 8 -> Pic with caption
    # """

    async def delete_all_slides():
        for i in range(len(root.slides) - 1, -1, -1):
            r_id = root.slides._sldIdLst[i].rId
            root.part.drop_rel(r_id)
            del root.slides._sldIdLst[i]

    async def create_title_slide(title, subtitle):
        layout = root.slide_layouts[0]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    async def create_section_header_slide(title):
        layout = root.slide_layouts[2]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title

    async def create_title_and_content_slide(title, content):
        layout = root.slide_layouts[1]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content

    async def create_title_and_content_and_image_slide(title, content, image_query):
        layout = root.slide_layouts[8]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[2].text = content

        try:
            image_data = await downloader.download(image_query, limit=1, adult_filter_off=True, timeout=15,
                                                   filter="+filterui:aspect-wide+filterui:imagesize-wallpaper+filterui:photo-photo")
            slide.placeholders[1].insert_picture(io.BytesIO(image_data))
        except Exception:
            pass

    async def find_text_in_between_tags(text, start_tag, end_tag):
        start_pos = text.find(start_tag)
        end_pos = text.find(end_tag)
        result = []
        while start_pos > -1 and end_pos > -1:
            text_between_tags = text[start_pos + len(start_tag):end_pos]
            result.append(text_between_tags)
            start_pos = text.find(start_tag, end_pos + len(end_tag))
            end_pos = text.find(end_tag, start_pos)
        res1 = "".join(result)
        res2 = re.sub(r"\[IMAGE\].*?\[/IMAGE\]", '', res1)
        if len(result) > 0:
            return res2
        else:
            return ""

    async def search_for_slide_type(text):
        tags = ["[L_TS]", "[L_CS]", "[L_IS]", "[L_THS]"]
        found_text = next((s for s in tags if s in text), None)
        return found_text

    async def parse_response(reply):
        list_of_slides = reply.split("[SLIDEBREAK]")
        for slide in list_of_slides:
            slide_type = await search_for_slide_type(slide)
            match slide_type:
                case ("[L_TS]"):
                    await create_title_slide(await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]"),
                                             await find_text_in_between_tags(str(slide), "[SUBTITLE]", "[/SUBTITLE]"))
                case ("[L_CS]"):
                    await create_title_and_content_slide("".join(await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                                                         "".join(await find_text_in_between_tags(str(slide), "[CONTENT]", "[/CONTENT]")))
                case ("[L_IS]"):
                    await create_title_and_content_and_image_slide("".join(await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                                                                   "".join(await find_text_in_between_tags(str(slide), "[CONTENT]", "[/CONTENT]")),
                                                                   "".join(await find_text_in_between_tags(str(slide), "[IMAGE]", "[/IMAGE]")))
                case ("[L_THS]"):
                    await create_section_header_slide("".join(await find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")))

    async def find_title():
        return root.slides[0].shapes.title.text

    await delete_all_slides()
    await parse_response(answer)
    buffer = io.BytesIO()
    root.save(buffer)
    pptx_bytes = buffer.getvalue()
    pptx_title = f"{await find_title()}.pptx"
    print(f"done {pptx_title}")

    return pptx_bytes, pptx_title
